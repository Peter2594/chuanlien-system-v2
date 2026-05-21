# -*- coding: utf-8 -*-
"""完成剩餘整併（PDF 已 ok，只處理 PPT 與 archive）"""
import os, shutil, sys

if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')

DOCS = "docs"
ARCHIVE = "docs/archive"
os.makedirs(ARCHIVE, exist_ok=True)


# ===== PPT 合併（穩定版）=====
# 改用「複製整份簡報」+ 「插入第二份的 slides」方法
print("[Step] Merge PPTs")
ppt1 = "docs/串連系統_產品介紹簡報.pptx"
ppt2 = "docs/串連系統_數字來源與解讀.pptx"
out_ppt = "docs/串連系統_產品簡報.pptx"

try:
    from pptx import Presentation
    from copy import deepcopy
    from pptx.util import Inches

    # 直接複製第一份作為基底（保留原始 layout 設定）
    shutil.copy(ppt1, out_ppt)
    print(f"  Base: copied {ppt1}")

    base = Presentation(out_ppt)
    add = Presentation(ppt2)
    print(f"  Adding {len(add.slides)} slides from {ppt2}")

    # 取 base 的第一個非空 layout 作為新 slide 的 layout
    blank_layout = None
    for lay in base.slide_layouts:
        blank_layout = lay
        break

    for src_slide in add.slides:
        new_slide = base.slides.add_slide(blank_layout)
        # 清掉新 slide 的預設 placeholder
        for ph in list(new_slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)
        # 複製 src_slide 所有 shape
        for shape in src_slide.shapes:
            el = shape.element
            new_el = deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    base.save(out_ppt)
    print(f"  -> {out_ppt} ({len(base.slides)} slides)")
except Exception as e:
    print(f"  PPT merge error: {e}")
    print(f"  Fallback: keeping {ppt1} as the final presentation")
    if os.path.exists(out_ppt):
        try:
            os.remove(out_ppt)
        except Exception:
            pass
    shutil.copy(ppt1, out_ppt)


# ===== 把所有舊檔搬到 archive =====
print()
print("[Step] Archive old files")
old_files = [
    "docs/串連系統_演算法說明.pdf",
    "docs/串連系統_演算法總覽.pdf",
    "docs/串連系統_完整技術手冊.pdf",
    "docs/串連系統_演算法視覺化手冊.pdf",
    "docs/串連系統_演算法深度解析.pdf",
    "docs/串連系統_答辯參數依據與標準答案.pdf",
    "docs/串連系統_產品介紹簡報.pptx",
    "docs/串連系統_數字來源與解讀.pptx",
    "docs/串連系統_公式設計推導.pdf",
]
for f in old_files:
    if os.path.exists(f):
        dst = os.path.join(ARCHIVE, os.path.basename(f))
        try:
            if os.path.exists(dst):
                os.remove(dst)
            shutil.move(f, dst)
            print(f"  -> archive/{os.path.basename(f)}")
        except PermissionError:
            print(f"  [LOCKED] {os.path.basename(f)} — 請手動關閉並再執行")

# ===== 結尾 =====
print()
print("=" * 60)
print("[DONE] Final files in docs/:")
print("=" * 60)
for f in sorted(os.listdir(DOCS)):
    full = os.path.join(DOCS, f)
    if os.path.isfile(full) and f.endswith((".pdf", ".pptx")):
        size = os.path.getsize(full) / 1024
        print(f"  {f}  ({size:.0f} KB)")
