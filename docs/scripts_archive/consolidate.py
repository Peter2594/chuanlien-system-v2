# -*- coding: utf-8 -*-
"""整併 9 → 4 份"""
import os, shutil, sys
from pypdf import PdfReader, PdfWriter

# 強制使用 UTF-8 輸出（避免 Windows cp950 編碼錯誤）
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')

DOCS = "docs"
ARCHIVE = "docs/archive"
os.makedirs(ARCHIVE, exist_ok=True)

# =====================
# 1. 合併 PDF
# =====================
print("=" * 60)
print("[Step 1] Merging PDFs")
print("=" * 60)

pdfs_to_merge = [
    ("docs/串連系統_完整技術手冊.pdf",    "Part 1: 完整技術手冊"),
    ("docs/串連系統_演算法深度解析.pdf",  "Part 2: 演算法深度解析"),
    ("docs/串連系統_公式設計推導.pdf",    "Part 3: 公式設計推導"),
]

writer = PdfWriter()
total_pages = 0
for path, name in pdfs_to_merge:
    if not os.path.exists(path):
        print(f"  [skip] {path}")
        continue
    reader = PdfReader(path)
    pages = len(reader.pages)
    print(f"  + {name} ({pages} pages)")
    for page in reader.pages:
        writer.add_page(page)
    writer.add_outline_item(name, total_pages)
    total_pages += pages

output_path = "docs/串連系統_完整技術文件.pdf"
with open(output_path, "wb") as f:
    writer.write(f)
print(f"  -> Output: {output_path} ({total_pages} pages total)")

# =====================
# 2. 合併 PPT
# =====================
print()
print("=" * 60)
print("[Step 2] Merging PPTs")
print("=" * 60)

from pptx import Presentation
from copy import deepcopy

ppts_to_merge = [
    "docs/串連系統_產品介紹簡報.pptx",
    "docs/串連系統_數字來源與解讀.pptx",
]


def merge_ppts(input_files, output_file):
    base = Presentation(input_files[0])
    print(f"  Base: {input_files[0]} ({len(base.slides)} slides)")
    for additional_path in input_files[1:]:
        if not os.path.exists(additional_path):
            continue
        additional = Presentation(additional_path)
        print(f"  + {additional_path} ({len(additional.slides)} slides)")
        for slide in additional.slides:
            slide_layout = base.slide_layouts[6]  # blank
            new_slide = base.slides.add_slide(slide_layout)
            for shape in slide.shapes:
                el = shape.element
                new_el = deepcopy(el)
                new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
    base.save(output_file)
    print(f"  -> Output: {output_file} ({len(base.slides)} slides)")


try:
    merge_ppts(ppts_to_merge, "docs/串連系統_產品簡報.pptx")
except Exception as e:
    print(f"  [PPT merge fallback]: {e}")
    if os.path.exists(ppts_to_merge[0]):
        shutil.copy(ppts_to_merge[0], "docs/串連系統_產品簡報.pptx")

# =====================
# 3. 重新命名保留的兩份
# =====================
print()
print("=" * 60)
print("[Step 3] Rename retained files")
print("=" * 60)

renames = [
    ("docs/串連系統_演算法視覺化手冊.pdf",       "docs/串連系統_視覺化手冊.pdf"),
    ("docs/串連系統_答辯參數依據與標準答案.pdf", "docs/串連系統_答辯速查.pdf"),
]
for old, new in renames:
    if os.path.exists(old):
        if os.path.exists(new):
            os.remove(new)
        shutil.copy(old, new)
        print(f"  + {old} -> {new}")

# =====================
# 4. 舊檔移至 archive/
# =====================
print()
print("=" * 60)
print("[Step 4] Archive old files")
print("=" * 60)

old_files = [
    "docs/串連系統_演算法說明.pdf",
    "docs/串連系統_演算法總覽.pdf",
    "docs/串連系統_完整技術手冊.pdf",
    "docs/串連系統_演算法視覺化手冊.pdf",
    "docs/串連系統_演算法深度解析.pdf",
    "docs/串連系統_答辯參數依據與標準答案.pdf",
    "docs/串連系統_產品介紹簡報.pptx",
    "docs/串連系統_數字來源與解讀.pptx",
    "docs/串連系統_公式設計推導.pdf",
]
for f in old_files:
    if os.path.exists(f):
        dst = os.path.join(ARCHIVE, os.path.basename(f))
        if os.path.exists(dst):
            os.remove(dst)
        shutil.move(f, dst)
        print(f"  -> archive/{os.path.basename(f)}")

# =====================
# 結尾
# =====================
print()
print("=" * 60)
print("[DONE] Final files in docs/:")
print("=" * 60)
for f in sorted(os.listdir(DOCS)):
    full = os.path.join(DOCS, f)
    if os.path.isfile(full) and f.endswith((".pdf", ".pptx")):
        size = os.path.getsize(full) / 1024
        print(f"  {f}  ({size:.0f} KB)")

print()
print("Backups in docs/archive/")
