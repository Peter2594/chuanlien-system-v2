# -*- coding: utf-8 -*-
"""收尾：從 archive 拿 PPT 合併 + 清理 docs/"""
import os, shutil, sys
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')

DOCS = "docs"
ARCHIVE = "docs/archive"

# ===== PPT 合併 =====
print("[Step] Merge PPTs (from archive)")
ppt1 = "docs/archive/串連系統_產品介紹簡報.pptx"
ppt2 = "docs/archive/串連系統_數字來源與解讀.pptx"
out_ppt = "docs/串連系統_產品簡報.pptx"

# 如果還有殘留的 .pptx 在 docs/，先搬走
for f in ["docs/串連系統_數字來源與解讀.pptx", "docs/串連系統_公式設計推導.pdf"]:
    if os.path.exists(f):
        dst = os.path.join(ARCHIVE, os.path.basename(f))
        if os.path.exists(dst):
            os.remove(f)  # archive 已有副本，直接刪 docs
        else:
            shutil.move(f, dst)
        print(f"  cleanup: {os.path.basename(f)}")

try:
    from pptx import Presentation
    from copy import deepcopy

    # 用 ppt1 作 base
    shutil.copy(ppt1, out_ppt)
    print(f"  Base: {ppt1}")
    base = Presentation(out_ppt)
    print(f"  Base has {len(base.slides)} slides, layouts={len(base.slide_layouts)}")

    add = Presentation(ppt2)
    print(f"  Adding {len(add.slides)} slides from {ppt2}")

    # 用第一個 layout（base 至少有 1 個）
    blank_layout = base.slide_layouts[0]

    for src_slide in add.slides:
        new_slide = base.slides.add_slide(blank_layout)
        # 清掉 placeholder
        for ph in list(new_slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)
        # 複製 shape
        for shape in src_slide.shapes:
            el = shape.element
            new_el = deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    base.save(out_ppt)
    print(f"  -> {out_ppt} ({len(base.slides)} slides)")
except Exception as e:
    print(f"  PPT merge error: {type(e).__name__}: {e}")
    if not os.path.exists(out_ppt):
        shutil.copy(ppt1, out_ppt)
        print(f"  Fallback: kept {ppt1} as final")

# ===== 結尾 =====
print()
print("=" * 60)
print("[DONE] Final files in docs/:")
print("=" * 60)
for f in sorted(os.listdir(DOCS)):
    full = os.path.join(DOCS, f)
    if os.path.isfile(full) and f.endswith((".pdf", ".pptx")):
        size = os.path.getsize(full) / 1024
        print(f"  {f}  ({size:.0f} KB)")

print()
print("Archive in docs/archive/")
for f in sorted(os.listdir(ARCHIVE)):
    full = os.path.join(ARCHIVE, f)
    if os.path.isfile(full):
        size = os.path.getsize(full) / 1024
        print(f"  {f}  ({size:.0f} KB)")
