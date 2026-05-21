# -*- coding: utf-8 -*-
"""串連系統 v2.2 — 簡報產生器 (Sequoia Memo 風)

設計準則：
- 純白底 + 深黑字 + 一個藍色 accent
- 中文細明體 + 英文 Times New Roman（襯線字）
- 大量留白、左對齊、無卡片、無色塊
- Roman numeral 章節編號
- 公式用 Times italic（學術數學風）
- 觀眾感受：在讀一份備忘錄而不是看 dashboard
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')

# ============ 配色 ============
BG     = RGBColor(0xFF, 0xFF, 0xFF)  # 純白
INK    = RGBColor(0x1A, 0x1A, 0x1A)  # 主文字（不純黑，柔和點）
MID    = RGBColor(0x44, 0x44, 0x44)  # 次文字
LITE   = RGBColor(0x88, 0x88, 0x88)  # 註解 / 頁碼
RULE   = RGBColor(0xCC, 0xCC, 0xCC)  # 分隔線
ACCENT = RGBColor(0x1E, 0x40, 0xAF)  # 深藍（key term 用）
QUOTE  = RGBColor(0x55, 0x55, 0x55)  # 引文色

# ============ 字型 ============
SERIF_CN = "PMingLiU"          # 細明體
SERIF_EN = "Times New Roman"
MONO     = "Consolas"
SANS_CN  = "Microsoft JhengHei"  # 只用在頁碼

# ============ 16:9 ============
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
M_LEFT  = Inches(1.4)   # 左邊界（給足夠呼吸）
M_RIGHT = Inches(1.4)
TEXT_W  = SLIDE_W - M_LEFT - M_RIGHT

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# ============ 通用元件 ============
def new_slide():
    s = prs.slides.add_slide(blank)
    # 白底
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    spTree = bg._element.getparent()
    spTree.remove(bg._element); spTree.insert(2, bg._element)
    return s


def add_text(slide, x, y, w, h, text, *,
             size=14, bold=False, italic=False, color=INK,
             font=SERIF_CN, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.55):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        # 同時設定 eastAsia 字型（影響中文渲染）
        from pptx.oxml.ns import qn
        rPr = run._r.get_or_add_rPr()
        rFonts = rPr.find(qn('a:rFonts'))
        if rFonts is None:
            rFonts = rPr.makeelement(qn('a:rFonts'), {})
            rPr.insert(0, rFonts)
        rFonts.set('eastAsia', SERIF_CN)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_rule(slide, y, *, width=None, color=RULE):
    """水平細線 — Sequoia Memo 標準分隔"""
    w = width or TEXT_W
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                M_LEFT, y, w, Emu(6350))
    ln.line.fill.background()
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    return ln


def add_page_footer(slide, num, total):
    """頁碼 + 組別小字 — 底部 sans"""
    add_text(slide, M_LEFT, Inches(7.05), Inches(8), Inches(0.3),
             "串連系統 v2.2", size=9, color=LITE, font=SANS_CN)
    add_text(slide, SLIDE_W - M_RIGHT - Inches(2), Inches(7.05),
             Inches(2), Inches(0.3),
             f"{num} / {total}",
             size=9, color=LITE, font=SANS_CN, align=PP_ALIGN.RIGHT)


def add_section_header(slide, num_roman, title, subtitle=None, y_start=Inches(1.0)):
    """Sequoia 標準頁頭：章節數字 / 標題 / 副標 / 分隔線"""
    # 章節 Roman numeral
    add_text(slide, M_LEFT, y_start, TEXT_W, Inches(0.35),
             num_roman, size=11, bold=True, color=ACCENT,
             font=SERIF_EN, line_spacing=1.0)
    # 標題
    add_text(slide, M_LEFT, y_start + Inches(0.4), TEXT_W, Inches(0.8),
             title, size=32, bold=True, color=INK, line_spacing=1.1)
    # 副標
    if subtitle:
        add_text(slide, M_LEFT, y_start + Inches(1.2), TEXT_W, Inches(0.5),
                 subtitle, size=15, italic=True, color=MID, line_spacing=1.3)
    # 分隔線
    rule_y = y_start + (Inches(1.85) if subtitle else Inches(1.35))
    add_rule(slide, rule_y)
    return rule_y + Inches(0.45)  # 回傳內文起始 y


# ============================================================
# Slide 1 — 封面
# ============================================================
def slide_cover():
    s = new_slide()
    # 上方標語（小字）
    add_text(s, M_LEFT, Inches(0.9), TEXT_W, Inches(0.4),
             "MANAGEMENT INFORMATION SYSTEMS · FINAL PROJECT",
             size=10, bold=True, color=LITE, font=SERIF_EN, line_spacing=1.0)
    add_text(s, M_LEFT, Inches(1.25), TEXT_W, Inches(0.4),
             "資訊管理導論　期末專案",
             size=12, color=LITE)

    # 主標題（最大字級）
    add_text(s, M_LEFT, Inches(2.8), TEXT_W, Inches(1.6),
             "串連系統",
             size=80, bold=True, color=INK, line_spacing=1.0)

    # 副標
    add_text(s, M_LEFT, Inches(4.2), TEXT_W, Inches(0.6),
             "管理決策支援平台",
             size=24, italic=True, color=MID, line_spacing=1.2)

    # 中段分隔線
    add_rule(s, Inches(5.3), width=Inches(1.2), color=ACCENT)

    # 作者與日期
    add_text(s, M_LEFT, Inches(5.5), TEXT_W, Inches(0.4),
             "第 13 組",
             size=13, bold=True, color=INK)
    add_text(s, M_LEFT, Inches(5.9), TEXT_W, Inches(0.4),
             "林聿平　組員 A　組員 B　組員 C　組員 D",
             size=12, color=MID)
    add_text(s, M_LEFT, Inches(6.35), TEXT_W, Inches(0.35),
             "2026 年 5 月",
             size=11, italic=True, color=LITE)


# ============================================================
# Slide 2 — Table of Contents
# ============================================================
def slide_agenda():
    s = new_slide()
    body_y = add_section_header(s, "CONTENTS", "目錄", None)

    sections = [
        ("I",    "問題",       "管理層看不見「正在發生什麼」"),
        ("II",   "解法",       "三個模組對應三個問題"),
        ("III",  "員工負載",   "把分散訊號加總為一個分數"),
        ("IV",   "組織健康度", "六維體檢看哪邊凹下去"),
        ("V",    "決策影響",   "扣掉大盤，看真實貢獻"),
        ("VI",   "方法論",     "整合管理層意見 → 等價關係 → 弱先驗 → 校準 → 敏感度"),
        ("VII",  "亮點功能",   "智能搜尋 · What-if · 部門網絡"),
        ("VIII", "競品分析",   "為什麼不用 Asana / Notion / Viva"),
        ("IX",   "效益與限制", "誠實揭露"),
    ]
    y = body_y
    for num, title, desc in sections:
        # Roman numeral
        add_text(s, M_LEFT, y, Inches(0.7), Inches(0.4),
                 num, size=12, bold=True, color=ACCENT,
                 font=SERIF_EN, line_spacing=1.0)
        # Title
        add_text(s, M_LEFT + Inches(0.7), y, Inches(3.0), Inches(0.4),
                 title, size=14, bold=True, color=INK, line_spacing=1.0)
        # Desc
        add_text(s, M_LEFT + Inches(3.7), y, Inches(6.8), Inches(0.4),
                 desc, size=12, italic=True, color=MID, line_spacing=1.0)
        y += Inches(0.42)

    add_page_footer(s, 2, 16)


# ============================================================
# Slide 3 — 問題
# ============================================================
def slide_problem():
    s = new_slide()
    body_y = add_section_header(
        s, "I.", "問題",
        "管理層每天面對的問題其實只有三個。")

    # 三段問題
    items = [
        ("誰快撐不住了？",
         "員工負載靠主管「感覺」，沒有客觀指標。"
         "資深主管報喜不報憂、新手主管把正常忙誤判成爆掉。"
         "每個主管的「忙」標準不同，導致資源無法跨部門公平分配。"),
        ("組織哪裡卡住了？",
         "跨部門卡點分散在週報、Email、Line 群組裡。"
         "管理層必須逐一追問才能拼出全貌，往往發現時已經拖了四週。"),
        ("我這個決策有效嗎？",
         "管理層每週做出大量決策，但事後幾乎沒有評估流程。"
         "下次遇到類似情況又憑直覺做一次，組織學不到任何東西。"),
    ]
    y = body_y
    for q, a in items:
        add_text(s, M_LEFT, y, TEXT_W, Inches(0.4),
                 q, size=18, bold=True, color=INK, line_spacing=1.2)
        add_text(s, M_LEFT, y + Inches(0.45), TEXT_W, Inches(1.0),
                 a, size=13, color=MID, line_spacing=1.6)
        y += Inches(1.3)

    add_page_footer(s, 3, 16)


# ============================================================
# Slide 4 — 解法
# ============================================================
def slide_solution():
    s = new_slide()
    body_y = add_section_header(
        s, "II.", "解法",
        "三個問題拆成三個模組，每個都有可解釋的數據產出。")

    mods = [
        ("員工負載分析",
         "誰快撐不住了",
         "把卡點、案件、交接、被提及四種訊號加權合計，配合時間衰減算分。"
         "產出是一張彩色長條圖，三秒看完全公司負載。",
         "經驗百分位 · 時間衰減 · Gini 不均度"),
        ("組織健康度雷達",
         "哪裡卡住了",
         "六個面向同時體檢（流動、卡點、決策、健康、負載、溝通），加權加總。"
         "產出是一張雷達圖加本週優先改善清單。",
         "六維加權評分 · 不對稱偵測"),
        ("決策影響評估",
         "決策有效嗎",
         "決策前後拍快照，扣掉「同期大盤本來會走到的位置」，得到相對表現。"
         "對應經濟學的 Difference-in-Differences 方法。",
         "同期校準 (DiD) · 線性回歸基準漂移"),
    ]
    y = body_y
    for name, q, desc, algo in mods:
        add_text(s, M_LEFT, y, Inches(4.0), Inches(0.35),
                 name, size=15, bold=True, color=INK, line_spacing=1.1)
        add_text(s, M_LEFT + Inches(4.0), y, Inches(5.5), Inches(0.35),
                 f"— {q}",
                 size=14, italic=True, color=ACCENT, line_spacing=1.1)
        add_text(s, M_LEFT, y + Inches(0.4), TEXT_W, Inches(0.7),
                 desc, size=12, color=MID, line_spacing=1.55)
        add_text(s, M_LEFT, y + Inches(1.15), TEXT_W, Inches(0.3),
                 f"演算法　{algo}",
                 size=10, italic=True, color=LITE, line_spacing=1.0)
        y += Inches(1.55)

    add_page_footer(s, 4, 16)


# ============================================================
# Slide 5 — 員工負載
# ============================================================
def slide_load():
    s = new_slide()
    body_y = add_section_header(
        s, "III.", "員工負載分析",
        "把分散在週報與交接中的工作訊號，加總為一個可比較的分數。")

    # 場景 quote
    add_text(s, M_LEFT, body_y, Inches(0.3), Inches(0.4),
             "「",
             size=24, bold=True, color=ACCENT, font=SERIF_EN, line_spacing=1.0)
    add_text(s, M_LEFT + Inches(0.3), body_y + Inches(0.1),
             TEXT_W - Inches(0.6), Inches(0.7),
             "我們組 17 個人，到底誰最操？我問了一圈，每個人都說自己很忙。",
             size=14, italic=True, color=QUOTE, line_spacing=1.4)
    add_text(s, M_LEFT + Inches(0.3), body_y + Inches(0.65),
             TEXT_W, Inches(0.35),
             "— 林經理",
             size=11, italic=True, color=LITE, line_spacing=1.0)

    # 公式
    formula_y = body_y + Inches(1.3)
    add_text(s, M_LEFT, formula_y, TEXT_W, Inches(0.3),
             "FORMULA",
             size=10, bold=True, color=ACCENT, font=SERIF_EN, line_spacing=1.0)
    add_text(s, M_LEFT, formula_y + Inches(0.35), TEXT_W, Inches(0.5),
             "Load  =  Σᵢ  wᵢ · Countᵢ · e^(−tᵢ/14)",
             size=20, italic=True, color=INK, font=SERIF_EN, line_spacing=1.2)
    add_text(s, M_LEFT, formula_y + Inches(0.9), TEXT_W, Inches(0.4),
             "每件事 × 重要性權重 × 時間衰減（半衰期 14 天，源自 Andy Grove 管理週期）",
             size=12, italic=True, color=MID, line_spacing=1.4)

    # 三層拆解
    add_rule(s, formula_y + Inches(1.5))
    layer_y = formula_y + Inches(1.7)

    layers = [
        ("時間衰減", "本週 1.0、1 週前 0.7、2 週前 0.4、3 週前 0.15、4 週後歸零。"),
        ("經驗百分位", "把絕對分數換成全公司排名百分比 — 林聿平 10.5 分對應 percentile 100%。"),
        ("Gini 不均度", "全員負載分佈的不平均程度，> 0.35 觸發「分配不均」警示。"),
    ]
    for name, desc in layers:
        add_text(s, M_LEFT, layer_y, Inches(2.5), Inches(0.3),
                 name, size=12, bold=True, color=INK, line_spacing=1.0)
        add_text(s, M_LEFT + Inches(2.5), layer_y, TEXT_W - Inches(2.5), Inches(0.6),
                 desc, size=11.5, color=MID, line_spacing=1.5)
        layer_y += Inches(0.55)

    add_page_footer(s, 5, 16)


# ============================================================
# Slide 6 — 健康度
# ============================================================
def slide_health():
    s = new_slide()
    body_y = add_section_header(
        s, "IV.", "組織健康度雷達",
        "把六個獨立指標加權成單一健康度分數，雷達圖視覺呈現。")

    add_text(s, M_LEFT, body_y, TEXT_W, Inches(1.4),
             "業績是落後指標，等財報出來再修就晚了。我們提前看『過程指標』——"
             "流動性、卡點密度、決策效率、成員健康、負載均衡、溝通對稱六個面向，"
             "用加權平均算出 0-100 的組織健康度分數。",
             size=13, color=MID, line_spacing=1.65)

    # 六維權重表
    add_text(s, M_LEFT, body_y + Inches(1.6), TEXT_W, Inches(0.3),
             "DIMENSIONS",
             size=10, bold=True, color=ACCENT, font=SERIF_EN, line_spacing=1.0)
    dims = [
        ("流動性",   "22%", "工作能否流轉的核心指標"),
        ("成員健康", "18%", "員工負載分佈"),
        ("卡點密度", "18%", "P95+ 卡點越少越健康"),
        ("溝通對稱", "15%", "部門互動雙向程度"),
        ("決策效率", "15%", "逾期決策越少越好"),
        ("負載均衡", "12%", "Gini + 過載人數複合警示"),
    ]
    y = body_y + Inches(2.0)
    for name, w, desc in dims:
        add_text(s, M_LEFT, y, Inches(2.5), Inches(0.3),
                 name, size=13, bold=True, color=INK, line_spacing=1.0)
        add_text(s, M_LEFT + Inches(2.5), y, Inches(1.0), Inches(0.3),
                 w, size=13, bold=True, color=ACCENT, font=SERIF_EN, line_spacing=1.0)
        add_text(s, M_LEFT + Inches(3.7), y, Inches(7.0), Inches(0.3),
                 desc, size=11.5, italic=True, color=MID, line_spacing=1.0)
        y += Inches(0.4)

    add_page_footer(s, 6, 16)


# ============================================================
# Slide 7 — 決策影響
# ============================================================
def slide_impact():
    s = new_slide()
    body_y = add_section_header(
        s, "V.", "決策影響評估",
        "扣掉「大盤本來會走到的位置」，得到該決策的真實貢獻。")

    # 場景
    add_text(s, M_LEFT, body_y, Inches(0.3), Inches(0.4),
             "「",
             size=24, bold=True, color=ACCENT, font=SERIF_EN, line_spacing=1.0)
    add_text(s, M_LEFT + Inches(0.3), body_y + Inches(0.1),
             TEXT_W - Inches(0.6), Inches(0.7),
             "我做了 5 個決策，系統說影響都是負的？可是這季大盤本來就跌啊。",
             size=14, italic=True, color=QUOTE, line_spacing=1.4)
    add_text(s, M_LEFT + Inches(0.3), body_y + Inches(0.65),
             TEXT_W, Inches(0.35),
             "— 部門主管",
             size=11, italic=True, color=LITE, line_spacing=1.0)

    # 公式
    formula_y = body_y + Inches(1.3)
    add_text(s, M_LEFT, formula_y, TEXT_W, Inches(0.3),
             "FORMULA · DIFFERENCE-IN-DIFFERENCES",
             size=10, bold=True, color=ACCENT, font=SERIF_EN, line_spacing=1.0)
    add_text(s, M_LEFT, formula_y + Inches(0.35), TEXT_W, Inches(0.5),
             "adjustedΔ  =  (after − before)  −  β · t",
             size=20, italic=True, color=INK, font=SERIF_EN, line_spacing=1.2)
    add_text(s, M_LEFT, formula_y + Inches(0.9), TEXT_W, Inches(0.4),
             "β 來自過去 12 週健康度的線性回歸斜率，t 是決策期間天數。",
             size=12, italic=True, color=MID, line_spacing=1.4)

    add_rule(s, formula_y + Inches(1.55))

    # 案例
    case_y = formula_y + Inches(1.75)
    add_text(s, M_LEFT, case_y, TEXT_W, Inches(0.3),
             "EXAMPLE",
             size=10, bold=True, color=ACCENT, font=SERIF_EN, line_spacing=1.0)
    add_text(s, M_LEFT, case_y + Inches(0.35), TEXT_W, Inches(1.5),
             "主管 A 健康度從 72 → 68，原始差 −4 看似糟糕。但同期大盤下滑 −13.5。"
             "adjustedΔ = −4 − (−13.5) = +9.5。大環境該讓你掉 13.5 分，"
             "你只掉 4 分，這 9.5 分的差距就是你的真實貢獻。",
             size=12.5, color=MID, line_spacing=1.6)

    add_page_footer(s, 7, 16)


# ============================================================
# Slide 8 — 方法論
# ============================================================
def slide_method():
    s = new_slide()
    body_y = add_section_header(
        s, "VI.", "方法論",
        "為什麼我們的人工參數不是憑空捏的 — 五階段反推校準。")

    steps = [
        ("整合管理層意見", "決定方向",
         "透過與管理層的訪談與業界文獻交叉驗證，整合出訊號重要性排序。"),
        ("等價關係", "決定初始尺度",
         "把抽象重要性翻成可比較的比例（如「一件卡點 ≈ 2.5 件案件」）。"),
        ("弱先驗", "不要一開始放大膽",
         "依排序給小差距（1.3 : 1.2 : 1.1 : 1.0），保留資料校準空間。"),
        ("測資校準", "決定調整幅度",
         "用 SEED 與極端測資跑排序，看哪些案例被低估／高估，反推調整。"),
        ("敏感度確認", "確認穩定性",
         "權重 ±20% 隨機擾動 500 次，Top-3 穩定 93.4%、Spearman ρ = 0.995。"),
    ]
    y = body_y
    for i, (name, role, desc) in enumerate(steps):
        # Roman numeral
        add_text(s, M_LEFT, y, Inches(0.6), Inches(0.4),
                 f"{i+1}.", size=14, bold=True, color=ACCENT,
                 font=SERIF_EN, line_spacing=1.0)
        # name + role
        add_text(s, M_LEFT + Inches(0.6), y, Inches(3.0), Inches(0.4),
                 name, size=14, bold=True, color=INK, line_spacing=1.0)
        add_text(s, M_LEFT + Inches(3.7), y, Inches(3.0), Inches(0.4),
                 f"— {role}",
                 size=13, italic=True, color=MID, line_spacing=1.0)
        # desc
        add_text(s, M_LEFT + Inches(0.6), y + Inches(0.4), TEXT_W - Inches(0.6),
                 Inches(0.6),
                 desc, size=11.5, color=MID, line_spacing=1.55)
        y += Inches(1.0)

    add_page_footer(s, 8, 16)


# ============================================================
# Slide 9 — BM25F 智能搜尋
# ============================================================
def slide_bm25f():
    s = new_slide()
    body_y = add_section_header(
        s, "VII.", "智能搜尋",
        "輸入「太洋」找得到「太洋證券法律意見書」— BM25F 加中文特化。")

    add_text(s, M_LEFT, body_y, TEXT_W, Inches(1.0),
             "搜尋的核心是 BM25F — 資訊檢索領域 50 年標準演算法，"
             "業界 Elasticsearch、Lucene、Notion 全用這套。"
             "我們在上面加了三層中文特化：n-gram 多粒度斷詞、子字串加成、同義詞合併。",
             size=13, color=MID, line_spacing=1.65)

    # 公式
    formula_y = body_y + Inches(1.2)
    add_text(s, M_LEFT, formula_y, TEXT_W, Inches(0.3),
             "FORMULA · BM25F (SIMPLIFIED)",
             size=10, bold=True, color=ACCENT, font=SERIF_EN, line_spacing=1.0)
    add_text(s, M_LEFT, formula_y + Inches(0.35), TEXT_W, Inches(0.5),
             "score  =  Σ  wᶠ · IDF · TF / (TF + k₁(1 − b + b · dl/avgdl))",
             size=16, italic=True, color=INK, font=SERIF_EN, line_spacing=1.2)
    add_text(s, M_LEFT, formula_y + Inches(0.9), TEXT_W, Inches(0.4),
             "欄位加權 × 詞稀有度 × 詞頻飽和 × 長度修正　·　k₁ = 1.5, b = 0.75 (Lucene 預設)",
             size=11.5, italic=True, color=MID, line_spacing=1.4)

    add_rule(s, formula_y + Inches(1.55))

    # 欄位權重
    weights_y = formula_y + Inches(1.75)
    add_text(s, M_LEFT, weights_y, TEXT_W, Inches(0.3),
             "FIELD WEIGHTS",
             size=10, bold=True, color=ACCENT, font=SERIF_EN, line_spacing=1.0)
    weight_text = "標題 5 · 標籤 4 · 摘要 2 · 結論 1.5 · 負責人 1 · 內文 1"
    add_text(s, M_LEFT, weights_y + Inches(0.35), TEXT_W, Inches(0.4),
             weight_text, size=13, color=INK, line_spacing=1.4)
    add_text(s, M_LEFT, weights_y + Inches(0.75), TEXT_W, Inches(0.4),
             "比例反映「資訊密度」— 標題濃縮整個案件意圖，留言通常離題。",
             size=11.5, italic=True, color=MID, line_spacing=1.4)

    add_page_footer(s, 9, 16)


# ============================================================
# Slide 10 — What-if
# ============================================================
def slide_whatif():
    s = new_slide()
    body_y = add_section_header(
        s, "VIII.", "What-if 模擬器",
        "在不影響實際資料的情況下，預先試算各種介入的健康度影響。")

    add_text(s, M_LEFT, body_y, TEXT_W, Inches(1.0),
             "管理層常面臨「不知道做了之後會不會更好」的決策困境。"
             "What-if 模擬器讓使用者調整四種介入，即時看到模擬後的健康度雷達圖對比，"
             "決定值不值得真的去執行。",
             size=13, color=MID, line_spacing=1.65)

    # 四種介入
    add_text(s, M_LEFT, body_y + Inches(1.2), TEXT_W, Inches(0.3),
             "FOUR LEVERS",
             size=10, bold=True, color=ACCENT, font=SERIF_EN, line_spacing=1.0)
    levers = [
        ("解卡點",   "勾選哪些卡點假設已解決"),
        ("加速決策", "勾選哪些逾期決策假設立刻完成"),
        ("簽收交接", "勾選哪些交接假設立刻簽收"),
        ("降低負載", "對特定員工拉滑桿減 0-80%"),
    ]
    y = body_y + Inches(1.65)
    for name, desc in levers:
        add_text(s, M_LEFT, y, Inches(2.5), Inches(0.35),
                 name, size=13, bold=True, color=INK, line_spacing=1.0)
        add_text(s, M_LEFT + Inches(2.5), y, TEXT_W - Inches(2.5), Inches(0.35),
                 desc, size=12, color=MID, line_spacing=1.0)
        y += Inches(0.45)

    add_rule(s, y + Inches(0.1))

    # 判定邏輯
    add_text(s, M_LEFT, y + Inches(0.3), TEXT_W, Inches(0.3),
             "VERDICT",
             size=10, bold=True, color=ACCENT, font=SERIF_EN, line_spacing=1.0)
    add_text(s, M_LEFT, y + Inches(0.65), TEXT_W, Inches(0.6),
             "Δ > 5 「強烈建議執行」 ·  0 ~ 5 「微改善」 ·  Δ < −5 「不建議」。"
             "技術上用 React 19 useDeferredValue 確保拖滑桿時不卡 UI。",
             size=12, italic=True, color=MID, line_spacing=1.5)

    add_page_footer(s, 10, 16)


# ============================================================
# Slide 11 — 部門網絡
# ============================================================
def slide_network():
    s = new_slide()
    body_y = add_section_header(
        s, "IX.", "部門互動網絡",
        "從圖論看「誰在單方面追著誰跑」— 不對稱互動即組織卡點訊號。")

    add_text(s, M_LEFT, body_y, TEXT_W, Inches(1.0),
             "健康的部門互動是雙向的（A 問 B 五次、B 也回 A 五次）。"
             "不健康的長這樣：A 問 B 九次、B 只回 A 一次 — A 在追、B 在躲。"
             "我們用統計閾值自動標出這類關係。",
             size=13, color=MID, line_spacing=1.65)

    # 公式
    formula_y = body_y + Inches(1.2)
    add_text(s, M_LEFT, formula_y, TEXT_W, Inches(0.3),
             "FORMULA",
             size=10, bold=True, color=ACCENT, font=SERIF_EN, line_spacing=1.0)
    add_text(s, M_LEFT, formula_y + Inches(0.35), TEXT_W, Inches(0.5),
             "Δ  =  (A → B)  −  (B → A)　·　閾值  =  1σ  ≈  5 次",
             size=18, italic=True, color=INK, font=SERIF_EN, line_spacing=1.2)
    add_text(s, M_LEFT, formula_y + Inches(0.9), TEXT_W, Inches(0.4),
             "經驗分佈標準差 σ ≈ 4.8，取 5 抓 top 16% 不平衡（對應 z-score > 1）。",
             size=12, italic=True, color=MID, line_spacing=1.4)

    add_rule(s, formula_y + Inches(1.55))

    # SEED 驗證
    add_text(s, M_LEFT, formula_y + Inches(1.75), TEXT_W, Inches(0.3),
             "VALIDATION",
             size=10, bold=True, color=ACCENT, font=SERIF_EN, line_spacing=1.0)
    add_text(s, M_LEFT, formula_y + Inches(2.1), TEXT_W, Inches(0.8),
             "主管心中有 3 對「卡」的部門，系統自動抓出 2 對（Recall 67%）。"
             "抓到的全部正確（Precision 100%）。學理對標 Burt《Structural Holes》(1992)。",
             size=12, color=MID, line_spacing=1.6)

    add_page_footer(s, 11, 16)


# ============================================================
# Slide 12 — 競品分析
# ============================================================
def slide_competitive():
    s = new_slide()
    body_y = add_section_header(
        s, "X.", "競品分析",
        "現成工具解任務，卻不解「決策」— 為什麼不用 Asana / Notion / Viva。")

    # Table
    headers = ["", "串連系統", "Notion", "Asana", "15Five", "MS Viva", "Excel + Line"]
    rows = [
        ("員工負載量化",   "●", "○", "◐", "◐", "●", "○"),
        ("跨部門卡點偵測", "●", "○", "◐", "○", "◐", "○"),
        ("決策事後追蹤",   "●", "○", "○", "○", "◐", "○"),
        ("組織健康度雷達", "●", "○", "○", "◐", "●", "○"),
        ("不對稱互動偵測", "●", "○", "○", "○", "○", "○"),
        ("導入成本",       "免費", "$", "$$", "$$$", "$$$", "免費"),
    ]
    col_w = [Inches(2.5), Inches(1.4), Inches(1.1), Inches(1.1),
             Inches(1.1), Inches(1.1), Inches(1.4)]
    tx = M_LEFT
    ty = body_y + Inches(0.1)

    # Header row
    x_acc = tx
    for ci, h in enumerate(headers):
        bold = (ci == 1)  # 自家欄位粗體
        col = ACCENT if ci == 1 else MID
        add_text(s, x_acc, ty, col_w[ci], Inches(0.4),
                 h, size=11.5, bold=bold, color=col,
                 align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER,
                 line_spacing=1.0)
        x_acc += col_w[ci]
    # Header rule
    add_rule(s, ty + Inches(0.45), color=INK)

    # Data rows
    row_y = ty + Inches(0.6)
    for ri, row in enumerate(rows):
        x_acc = tx
        for ci, val in enumerate(row):
            is_self = (ci == 1)
            col = ACCENT if is_self else (INK if ci == 0 else MID)
            bold = is_self or (ci == 0)
            font_use = SERIF_CN if ci == 0 else SERIF_EN
            add_text(s, x_acc, row_y, col_w[ci], Inches(0.4),
                     val, size=12.5, bold=bold, color=col, font=font_use,
                     align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER,
                     line_spacing=1.0)
            x_acc += col_w[ci]
        # 細分隔線
        if ri < len(rows) - 1:
            add_rule(s, row_y + Inches(0.42), color=RULE)
        row_y += Inches(0.5)

    # legend
    add_text(s, M_LEFT, row_y + Inches(0.2), TEXT_W, Inches(0.3),
             "●  完整支援　 ◐  部分支援　 ○  不支援",
             size=10, italic=True, color=LITE)

    add_page_footer(s, 12, 16)


# ============================================================
# Slide 13 — Tech Stack
# ============================================================
def slide_tech():
    s = new_slide()
    body_y = add_section_header(
        s, "XI.", "技術架構",
        "前端即時 · 後端輕量 · 雲端託管。")

    layers = [
        ("Frontend",   "React 19  ·  TypeScript  ·  Vite 6",
                       "Tailwind v4  ·  Recharts  ·  Framer Motion"),
        ("Data Layer", "Firebase Auth  ·  Cloud Firestore",
                       "即時訂閱 · 離線同步"),
        ("Algorithms", "BM25F  ·  Gini  ·  Empirical Percentile",
                       "Cohort Regression  ·  全部純前端計算"),
        ("Deployment", "Vercel  ·  GitHub Actions",
                       "CI / CD 自動部署"),
    ]
    y = body_y
    for name, line1, line2 in layers:
        add_text(s, M_LEFT, y, Inches(2.0), Inches(0.3),
                 name, size=12, bold=True, color=ACCENT,
                 font=SERIF_EN, line_spacing=1.0)
        add_text(s, M_LEFT + Inches(2.2), y, TEXT_W - Inches(2.2), Inches(0.3),
                 line1, size=12.5, color=INK, line_spacing=1.0)
        add_text(s, M_LEFT + Inches(2.2), y + Inches(0.35), TEXT_W - Inches(2.2),
                 Inches(0.3),
                 line2, size=11, italic=True, color=MID, line_spacing=1.0)
        y += Inches(0.95)

    add_page_footer(s, 13, 16)


# ============================================================
# Slide 14 — 預期效益
# ============================================================
def slide_outcome():
    s = new_slide()
    body_y = add_section_header(
        s, "XII.", "預期效益",
        "從「主管直覺」到「數據共識」。")

    outcomes = [
        ("90 秒",
         "員工負載狀態",
         "原本：跟 17 人輪流訪談 30 分鐘　·　現在：打開儀表板 90 秒掃完。"),
        ("67%",
         "自動抓到卡點",
         "原本：主管心中的卡點靠記憶　·　現在：系統幫你抓出 2/3。"),
        ("100%",
         "決策被回看",
         "原本：做完就忘、下次照樣憑感覺　·　現在：每個決策強制事後評估。"),
        ("0 元",
         "授權成本",
         "全用開源 + Firebase 免費額度，中型公司可零成本試用。"),
    ]
    y = body_y
    for big, label, desc in outcomes:
        # 大數字
        add_text(s, M_LEFT, y, Inches(2.0), Inches(0.7),
                 big, size=36, bold=True, color=ACCENT,
                 font=SERIF_EN, line_spacing=1.0)
        # 標籤
        add_text(s, M_LEFT + Inches(2.2), y + Inches(0.08), TEXT_W - Inches(2.2),
                 Inches(0.35),
                 label, size=14, bold=True, color=INK, line_spacing=1.0)
        # 描述
        add_text(s, M_LEFT + Inches(2.2), y + Inches(0.45), TEXT_W - Inches(2.2),
                 Inches(0.4),
                 desc, size=11.5, italic=True, color=MID, line_spacing=1.4)
        y += Inches(1.05)

    add_page_footer(s, 14, 16)


# ============================================================
# Slide 15 — 限制
# ============================================================
def slide_limits():
    s = new_slide()
    body_y = add_section_header(
        s, "XIII.", "已知限制",
        "誠實揭露 — 我們知道哪裡還不夠好。")

    items = [
        ("SEED 並非真實資料",
         "所有參數都在團隊模擬資料上校準。下個階段需取得真實公司資料重跑反推。"),
        ("無外部 ground truth",
         "週報品質、搜尋準確度尚未跟真實主管標註比對，目前僅做內部一致性測試。"),
        ("樣本量級偏小",
         "17 員工的 percentile 與 Gini 有統計雜訊。保留是為了未來放大至 100+ 人。"),
        ("週報品質用字數當代理",
         "無法分辨「精準 80 字」與「灌水 80 字」。已有有寫卡點欄位的補救機制，但只能擋懶人型灌水。"),
    ]
    y = body_y
    for title, desc in items:
        add_text(s, M_LEFT, y, TEXT_W, Inches(0.35),
                 f"— {title}",
                 size=14, bold=True, color=INK, line_spacing=1.1)
        add_text(s, M_LEFT + Inches(0.3), y + Inches(0.4), TEXT_W - Inches(0.3),
                 Inches(0.7),
                 desc, size=12, color=MID, line_spacing=1.55)
        y += Inches(1.15)

    add_page_footer(s, 15, 16)


# ============================================================
# Slide 16 — 結語
# ============================================================
def slide_closing():
    s = new_slide()

    # 章節編號
    add_text(s, M_LEFT, Inches(1.0), TEXT_W, Inches(0.35),
             "XIV.",
             size=11, bold=True, color=ACCENT, font=SERIF_EN, line_spacing=1.0)
    add_text(s, M_LEFT, Inches(1.4), TEXT_W, Inches(0.6),
             "結語",
             size=24, bold=True, color=INK, line_spacing=1.1)
    add_rule(s, Inches(2.25))

    # 主 take-away（大字置中）
    add_text(s, M_LEFT, Inches(3.3), TEXT_W, Inches(1.0),
             "工程現實 ≠ 學術理想",
             size=56, bold=True, color=INK,
             align=PP_ALIGN.CENTER, line_spacing=1.1)

    add_rule(s, Inches(4.6), width=Inches(1.2), color=ACCENT)

    add_text(s, M_LEFT, Inches(4.9), TEXT_W, Inches(0.8),
             "比起完美的公式，我們更相信可以解釋的數字。",
             size=18, italic=True, color=MID,
             align=PP_ALIGN.CENTER, line_spacing=1.6)

    # 底部署名
    add_text(s, M_LEFT, Inches(6.6), TEXT_W, Inches(0.3),
             "串連系統 v2.2  ·  資管導論 第 13 組  ·  2026 年 5 月",
             size=10, italic=True, color=LITE,
             align=PP_ALIGN.CENTER, font=SERIF_EN)


# ============================================================
slide_cover()
slide_agenda()
slide_problem()
slide_solution()
slide_load()
slide_health()
slide_impact()
slide_method()
slide_bm25f()
slide_whatif()
slide_network()
slide_competitive()
slide_tech()
slide_outcome()
slide_limits()
slide_closing()

out = "docs/串連系統_簡報.pptx"
prs.save(out)
print(f"OK -> {out} ({len(prs.slides)} slides)")
